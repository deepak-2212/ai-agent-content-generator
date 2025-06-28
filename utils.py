from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation

def save_as_pdf(content, filename):
    c = canvas.Canvas(filename, pagesize=letter)
    width, height = letter
    lines = content.split('\n')
    y = height - 40
    for line in lines:
        c.drawString(40, y, line)
        y -= 15
        if y < 40:
            c.showPage()
            y = height - 40
    c.save()

def save_as_pptx(content, filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.shapes.placeholders[1]
    title.text = "Generated Content"
    body.text = content
    prs.save(filename)
