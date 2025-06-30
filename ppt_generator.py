from pptx import Presentation
from pptx.util import Inches
import requests
from bs4 import BeautifulSoup

def scrape_image(topic):
    query = "+".join(topic.split())
    url = f"https://www.google.com/search?q={query}&tbm=isch"
    headers = {"User-Agent": "Mozilla/5.0"}
    response = requests.get(url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    img = soup.find("img")
    if img and img.get("src"):
        return img["src"]
    return None

def generate_ppt(topic, content, output_path="output/output.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = topic

    img_url = scrape_image(topic)
    if img_url:
        img_data = requests.get(img_url).content
        with open("output/temp.jpg", "wb") as f:
            f.write(img_data)
        slide.shapes.add_picture("output/temp.jpg", Inches(1), Inches(2), Inches(4), Inches(3))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1]
    body.text = content

    prs.save(output_path)
