from fastapi import FastAPI, Form
from fastapi.responses import HTMLResponse, FileResponse
from generator import generate_content
from fpdf import FPDF
from pptx import Presentation
import os

app = FastAPI()

@app.get("/", response_class=HTMLResponse)
async def root():
    return """
    <html>
        <head><title>Content Generator AI</title></head>
        <body>
            <h1>✅ Content Generator AI</h1>
            <form action="/generate" method="post">
                <input type="text" name="prompt" placeholder="Enter your prompt" required/>
                <select name="type">
                    <option value="blog">Blog (PDF)</option>
                    <option value="ppt">PPT</option>
                    <option value="report">Report (PDF)</option>
                </select>
                <button type="submit">Generate</button>
            </form>
        </body>
    </html>
    """

@app.post("/generate")
async def generate(prompt: str = Form(...), type: str = Form(...)):
    content = generate_content(prompt, type)

    if type == "ppt":
        filename = "output.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content

        # Split content into slides by 200 chars for example
        chunks = [content[i:i+200] for i in range(0, len(content), 200)]
        for i, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Slide {i+1}"
            slide.shapes.placeholders[1].text = chunk

        prs.save(filename)
        return FileResponse(path=filename, filename=filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    else:
        filename = "output.pdf"
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, content)
        pdf.output(filename)
        return FileResponse(path=filename, filename=filename, media_type='application/pdf')
